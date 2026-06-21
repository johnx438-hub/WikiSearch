"""
PPTX 转换器 — 将 PowerPoint 演示文稿转换为结构化 Markdown

核心功能:
1. 提取每张幻灯片的标题、正文、备注
2. 按幻灯片生成独立 MD 文件
3. 生成 _index.md 导航页
4. YAML Frontmatter 元数据注入

使用示例:
    converter = PptxConverter()
    result = converter.convert("/path/to/presentation.pptx")
"""

import os
from pathlib import Path
from typing import List, Dict, Optional
from dataclasses import dataclass, field


try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


# ==================== 数据结构 ====================

@dataclass
class SlideResult:
    """单张幻灯片转换结果"""
    title: str                    # 幻灯片标题
    content: str                  # Markdown 正文内容
    slide_number: int             # 幻灯片序号（1-indexed）
    notes: str = ""               # 演讲者备注
    metadata: Dict = field(default_factory=dict)  # YAML Frontmatter


# ==================== PPTX 转换器 ====================

class PptxConverter:
    """
    PowerPoint (.pptx) → Markdown 转换器。
    
    Attributes:
        extract_images: 是否提取图片到 assets/（默认 False，Phase 2实现）
    """
    
    def __init__(self, extract_images: bool = False):
        self.extract_images = extract_images
    
    def convert(self, filepath: str, output_dir: Optional[str] = None) -> Dict:
        """
        将 PPTX 文件转换为 Markdown。
        
        Args:
            filepath: .pptx 文件路径
            output_dir: 输出目录（如果提供，会生成树状 MD 文件）
            
        Returns:
            {
                'sections': [SlideResult, ...],
                'total_slides': int,
                'has_notes': bool,
            }
        """
        filepath = Path(filepath)
        if not filepath.exists():
            raise FileNotFoundError(f"文件不存在: {filepath}")
        
        # 1. 加载演示文稿
        prs = Presentation(str(filepath))
        
        # 2. 提取每张幻灯片
        slides = self._extract_slides(prs)
        
        # 3. 生成树状 MD 文件（如果提供了输出目录）
        if output_dir:
            self._write_tree_structure(
                sections=slides,
                base_dir=output_dir,
                doc_name=filepath.stem
            )
        
        return {
            'sections': slides,
            'total_slides': len(prs.slides),
            'has_notes': any(s.notes for s in slides),
        }
    
    def _extract_slides(self, prs: Presentation) -> List[SlideResult]:
        """
        提取所有幻灯片内容。
        
        Args:
            prs: python-pptx Presentation 对象
            
        Returns:
            [SlideResult, ...]
        """
        slides = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            # 提取标题
            title = self._extract_title(slide)
            
            # 提取正文内容（ bullet points ）
            content = self._extract_content(slide)
            
            # 提取备注
            notes = self._extract_notes(slide)
            
            # 构建元数据
            metadata = {
                'source': 'pptx',
                'slide_number': idx,
                'total_slides': len(prs.slides),
                'layout': slide.slide_layout.name if hasattr(slide, 'slide_layout') else 'Unknown',
            }
            
            slides.append(SlideResult(
                title=title or f"Slide {idx}",
                content=content,
                slide_number=idx,
                notes=notes,
                metadata=metadata
            ))
        
        return slides
    
    def _extract_title(self, slide) -> Optional[str]:
        """
        提取幻灯片标题。
        
        Args:
            slide: python-pptx Slide 对象
            
        Returns:
            标题文本或 None
        """
        # 尝试获取标题形状
        for shape in slide.shapes:
            if hasattr(shape, "name") and "Title" in shape.name:
                return shape.text_frame.text.strip()
        
        # 回退：取第一个非空文本
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    return text[:100]  # 限制长度
        
        return None
    
    def _extract_content(self, slide) -> str:
        """
        提取幻灯片正文内容。
        
        Args:
            slide: python-pptx Slide 对象
            
        Returns:
            Markdown 格式的内容
        """
        lines = []
        seen_texts = set()  # 避免重复
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # 跳过标题（已单独提取）
            if hasattr(shape, "name") and "Title" in shape.name:
                continue
            
            text_frame = shape.text_frame
            
            for para in text_frame.paragraphs:
                # 提取段落文本和格式信息
                para_text = self._format_paragraph(para)
                
                if para_text and para_text not in seen_texts:
                    seen_texts.add(para_text)
                    
                    # 判断是否为列表项（简单启发式）
                    text_stripped = para_text.strip()
                    if any(text_stripped.startswith(prefix) 
                           for prefix in ['•', '-', '*', '1.', '2.', '3.']):
                        lines.append(f"- {text_stripped}")
                    elif len(text_stripped) > 50:  # 长文本当段落
                        lines.append(f"\n{text_stripped}\n")
                    else:  # 短文本当列表项
                        lines.append(f"- {text_stripped}")
        
        return "\n".join(lines) if lines else ""
    
    def _format_paragraph(self, para) -> str:
        """
        格式化段落（保留加粗/斜体）。
        
        Args:
            para: python-pptx Paragraph 对象
            
        Returns:
            Markdown 格式文本
        """
        parts = []
        for run in para.runs:
            text = run.text.strip()
            if not text:
                continue
            
            # 基础格式标记
            if run.font.bold and run.font.italic:
                text = f"***{text}***"
            elif run.font.bold:
                text = f"**{text}**"
            elif run.font.italic:
                text = f"_{text}_"
            
            parts.append(text)
        
        return "".join(parts)
    
    def _extract_notes(self, slide) -> str:
        """
        提取演讲者备注。
        
        Args:
            slide: python-pptx Slide 对象
            
        Returns:
            备注文本
        """
        try:
            notes_slide = slide.notes_slide
            if notes_slide:
                notes_shape = notes_slide.notes_text_frame
                return notes_shape.text.strip()
        except (AttributeError, IndexError):
            pass
        
        return ""
    
    def _write_tree_structure(self, sections: List[SlideResult], 
                               base_dir: str, doc_name: str):
        """
        将转换结果写入树状 MD 文件结构。
        
        Args:
            sections: 幻灯片列表
            base_dir: 基础输出目录
            doc_name: 文档名称（不含扩展名）
        """
        # 创建文档根目录
        doc_dir = os.path.join(base_dir, doc_name)
        os.makedirs(doc_dir, exist_ok=True)
        
        # 生成 _index.md（导航页）
        self._write_index_md(sections=sections, 
                            index_path=os.path.join(doc_dir, "_index.md"))
        
        # 按幻灯片生成独立 MD 文件
        for slide in sections:
            filename = f"Slide_{slide.slide_number:02d}_{self._sanitize_filename(slide.title)}.md"
            filepath = os.path.join(doc_dir, filename)
            
            self._write_slide_md(
                section=slide,
                doc_name=doc_name,
                output_path=filepath
            )
    
    def _write_index_md(self, sections: List[SlideResult], index_path: str):
        """
        生成 _index.md（幻灯片导航）。
        
        Args:
            sections: 幻灯片列表
            index_path: 输出路径
        """
        lines = [
            f"# {sections[0].metadata.get('source', 'PPTX').upper()} 演示文稿",
            "",
            "## 📑 幻灯片导航",
            ""
        ]
        
        for slide in sections:
            filename = f"Slide_{slide.slide_number:02d}_{self._sanitize_filename(slide.title)}.md"
            lines.append(f"- [{slide.slide_number}. {slide.title}]({filename})")
        
        # 添加统计信息
        lines.extend([
            "",
            "## 📊 文档统计",
            f"- **总幻灯片数**: {len(sections)}",
            f"- **含备注**: {'是' if any(s.notes for s in sections) else '否'}",
        ])
        
        with open(index_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
    
    def _write_slide_md(self, section: SlideResult, doc_name: str, output_path: str):
        """
        生成单张幻灯片的 MD 文件。
        
        Args:
            section: 幻灯片结果
            doc_name: 文档名称
            output_path: 输出路径
        """
        # YAML Frontmatter
        frontmatter = [
            "---",
            f"source: \"{doc_name}.pptx\"",
            f"slide_number: {section.slide_number}",
            f"title: \"{section.title}\"",
            f"layout: \"{section.metadata.get('layout', 'Unknown')}\"",
            "---",
            ""
        ]
        
        # 标题 + 内容
        content_lines = [
            '\n'.join(frontmatter),
            f"# {section.title}",
            "",
            section.content,
        ]
        
        # 添加备注（如果有）
        if section.notes:
            content_lines.extend([
                "",
                "---",
                f"> 📝 **演讲者备注**:",
                f"> {section.notes}",
            ])
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(content_lines))
    
    def _sanitize_filename(self, text: str) -> str:
        """
        清理文件名中的非法字符。
        
        Args:
            text: 原始文本
            
        Returns:
            安全的文件名
        """
        import re
        sanitized = re.sub(r'[\\/:*?"<>|]', '_', text)
        if len(sanitized) > 30:
            sanitized = sanitized[:27] + '...'
        return sanitized.strip()
